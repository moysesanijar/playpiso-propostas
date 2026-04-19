from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from services.config import OUTPUT_DIR, TEMPLATE_BY_TYPE
from services.models import ProposalData, QuadraItem
from services.text_rules import (
    build_alambrado_text,
    build_iluminacao_text,
    build_materiais_pedreira_text,
    build_piso_text,
    build_playcushion_text,
    build_project_summary,
    build_responsabilidades_contratada,
    build_responsabilidades_contratante,
    format_area_pt,
    format_decimal_pt,
)


def _sanitize_filename(value: str) -> str:
    value = value.strip()
    value = re.sub(r"[^A-Za-zÀ-ÿ0-9._-]+", "_", value)
    return value.strip("_") or "proposta"


def _format_date_br(value) -> str:
    return value.strftime("%d/%m/%Y")


def _replace_in_run(run, replacements: dict[str, str]) -> bool:
    """Replace placeholders in a single run's text, preserving all formatting.
    Returns True if any replacement was made."""
    original = run.text
    new_text = original
    for source, target in replacements.items():
        new_text = new_text.replace(source, target)
    if new_text != original:
        run.text = new_text
        return True
    return False


def _replace_in_paragraph(paragraph, replacements: dict[str, str]) -> None:
    """Replace placeholders across all runs in a paragraph.

    Handles the common case where a placeholder like {{foo}} is split across
    multiple runs. We first check if the full paragraph text contains any
    placeholder; if it does but no individual run does, we consolidate the
    paragraph text into the first run (preserving that run's formatting) and
    clear the rest.
    """
    # Fast path: try replacing in individual runs first
    for run in paragraph.runs:
        _replace_in_run(run, replacements)

    # Check if any placeholder still spans multiple runs
    full_text = "".join(r.text for r in paragraph.runs)
    for placeholder in replacements:
        if placeholder in full_text:
            # Placeholder still present after per-run replacement →
            # it must be split. Consolidate into first run.
            new_text = full_text
            for source, target in replacements.items():
                new_text = new_text.replace(source, target)
            if paragraph.runs:
                paragraph.runs[0].text = new_text
                for run in paragraph.runs[1:]:
                    run.text = ""
            break


def _replace_in_text_frame(text_frame, replacements: dict[str, str]) -> None:
    for paragraph in text_frame.paragraphs:
        _replace_in_paragraph(paragraph, replacements)


def _replace_everywhere(prs: Presentation, replacements: dict[str, str]) -> None:
    """Walk every shape on every slide and replace placeholders in-place,
    preserving all XML formatting attributes."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                _replace_in_text_frame(shape.text_frame, replacements)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        _replace_in_text_frame(cell.text_frame, replacements)


def _get_slide_index_by_content(prs: Presentation, marker: str) -> int | None:
    """Return the index (0-based) of the first slide containing marker text."""
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                if marker in shape.text_frame.text:
                    return i
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if marker in cell.text_frame.text:
                            return i
    return None


def _hide_slide_by_index(prs: Presentation, slide_index_zero_based: int) -> None:
    if 0 <= slide_index_zero_based < len(prs.slides):
        slide_id_list = prs.slides._sldIdLst
        slide_id_list.remove(slide_id_list[slide_index_zero_based])


def _build_replacements(data: ProposalData, item: QuadraItem) -> dict[str, str]:
    # Alambrado area (perimeter * height logic: comp+larg both sides)
    alambrado_area = format_area_pt(
        (item.alambrado_comprimento * 2 + item.alambrado_largura * 2) * item.alambrado_altura
    )
    # Iluminação quantity = número de postes (conjunto)
    qtd_iluminacao = str(item.quantidade_postes)
    # Playcushion area = quadra area (only shown when Cushion type)
    qtd_playcushion = format_area_pt(item.area_total) if "Cushion" in item.tipo_quadra else "0"

    return {
        # ── Dados da proposta ────────────────────────────────────────────
        "{{numero_proposta}}": data.numero_proposta,
        "{{data_solicitacao}}": _format_date_br(data.data_solicitacao),
        "{{data_envio}}": _format_date_br(data.data_envio),
        # ── Dados do cliente ─────────────────────────────────────────────
        "{{cliente_nome}}": data.nome_cliente,
        "{{cliente_nome_cnpj}}": data.cliente_nome_cnpj,
        "{{cliente_contato}}": data.contato,
        "{{cliente_telefone}}": data.telefone,
        "{{cliente_email}}": data.email,
        "{{cliente_endereco}}": data.endereco,
        "{{local_obra}}": data.local_obra_completo,
        # ── Sumário ──────────────────────────────────────────────────────
        "{{resumo_projeto}}": build_project_summary(data),
        # ── Quadra / item ────────────────────────────────────────────────
        "{{tipo_quadra}}": item.tipo_quadra,
        "{{quantidade}}": str(item.quantidade),
        "{{comprimento}}": format_decimal_pt(item.comprimento, 2),
        "{{largura}}": format_decimal_pt(item.largura, 2),
        "{{area_total}}": format_area_pt(item.area_total),
        # ── Slides de conteúdo ───────────────────────────────────────────
        "{{texto_piso}}": build_piso_text(item),
        "{{texto_playcushion}}": build_playcushion_text(item),
        "{{texto_alambrado}}": build_alambrado_text(item),
        "{{texto_iluminacao}}": build_iluminacao_text(item),
        "{{texto_responsabilidades_contratada}}": build_responsabilidades_contratada(item),
        "{{texto_responsabilidades_contratante}}": build_responsabilidades_contratante(item),
        "{{texto_materiais_pedreira}}": build_materiais_pedreira_text(item),
        # ── Tabela de investimento ───────────────────────────────────────
        "{{tipo_estrutura}}": item.tipo_estrutura,
        "{{galvanizacao}}": item.galvanizacao,
        "{{travamento_tipo}}": item.travamento_tipo,
        "{{qtd_alambrado}}": alambrado_area,
        "{{qtd_iluminacao}}": qtd_iluminacao,
        "{{qtd_playcushion}}": qtd_playcushion,
        # Placeholders individuais na linha de iluminação da tabela de investimento
        "{{iluminacao_quantidade}}": str(item.iluminacao_quantidade),
        "{{potencia}}": str(item.potencia),
        "{{quantidade_postes}}": str(item.quantidade_postes),
        "{{altura_poste}}": format_decimal_pt(item.altura_poste, 1),
        # ── Tabela de materiais de pedreira (slide 14) ───────────────────
        "{{qtd_blocos}}": str(item.qtd_blocos),
        "{{qtd_areia_media}}": format_decimal_pt(item.qtd_areia_media, 2),
        "{{qtd_brita_1}}": format_decimal_pt(item.qtd_brita_1, 2),
        "{{qtd_po_de_pedra}}": format_decimal_pt(item.qtd_po_de_pedra, 2),
        "{{qtd_cimento}}": str(item.qtd_cimento),
        "{{qtd_lona_plastica}}": format_area_pt(item.qtd_lona_plastica),
        "{{qtd_pedrisco_limpo}}": format_decimal_pt(item.qtd_pedrisco_limpo, 2),
    }


def generate_ppt(data: ProposalData) -> Path:
    if not data.items:
        raise ValueError("A proposta não possui itens.")

    first_item = data.items[0]

    if first_item.tipo_quadra not in TEMPLATE_BY_TYPE:
        raise ValueError(f"Tipo de quadra sem template mapeado: {first_item.tipo_quadra}")

    template_path = Path(TEMPLATE_BY_TYPE[first_item.tipo_quadra])
    if not template_path.exists():
        raise FileNotFoundError(f"Template PPT não encontrado em: {template_path}")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(template_path))

    replacements = _build_replacements(data, first_item)
    _replace_everywhere(prs, replacements)

    # Ocultar slide do Playcushion (índice 9 = slide 10) se não for tipo Cushion
    if "Cushion" not in first_item.tipo_quadra:
        _hide_slide_by_index(prs, 9)

    # Ocultar slide de materiais de pedreira (último slide) se material for do cliente
    if first_item.material_pedreira.strip().lower() == "cliente":
        pedreira_idx = _get_slide_index_by_content(prs, "{{qtd_blocos}}")
        if pedreira_idx is None:
            # Placeholder já foi substituído — buscar pelo texto da tabela
            pedreira_idx = _get_slide_index_by_content(prs, "Areia Média")
        if pedreira_idx is not None:
            _hide_slide_by_index(prs, pedreira_idx)

    file_name = (
        f"{_sanitize_filename(data.numero_proposta)}_"
        f"{_sanitize_filename(data.nome_cliente)}.pptx"
    )
    output_path = OUTPUT_DIR / file_name
    prs.save(str(output_path))
    return output_path